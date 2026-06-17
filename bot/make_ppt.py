#!/usr/bin/env python3
# 用模板生成本周组会 PPT：把模板里的占位符 {{DATE}} / {{TITLE}} 替换成本周四的日期。
# 依赖： pip install python-pptx
#
# 模板放在 bot/template.pptx，在需要填日期的文本框里写占位符，例如标题页写：{{TITLE}}
# 占位符： {{DATE}} -> 2026-06-25    {{TITLE}} -> 2026.06.25-checkin
#
# 用法： python3 bot/make_ppt.py
# 输出： out/<YYYY-MM-DD>-checkin.pptx，并把 文件路径/日期/标题 打印为 GITHUB_OUTPUT 友好格式

import os
import sys
from datetime import datetime, timedelta, timezone

try:
    from pptx import Presentation
except ImportError:
    print("缺少 python-pptx，请先： pip install python-pptx", file=sys.stderr)
    sys.exit(1)

TEMPLATE = os.path.join(os.path.dirname(__file__), "template.pptx")


def this_thursday():
    """本周四（Asia/Shanghai）。今天是周四就取今天，否则取本周内的周四。"""
    now = datetime.now(timezone(timedelta(hours=8)))
    # weekday(): 周一=0 ... 周四=3
    offset = 3 - now.weekday()
    thu = now + timedelta(days=offset)
    return thu.date()


def replace_placeholders(prs, mapping):
    def fix_text_frame(tf):
        for para in tf.paragraphs:
            # 先合并整段文字做替换，再写回第一个 run，避免占位符被拆散到多个 run
            full = "".join(run.text for run in para.runs)
            new = full
            for k, v in mapping.items():
                new = new.replace(k, v)
            if new != full and para.runs:
                para.runs[0].text = new
                for run in para.runs[1:]:
                    run.text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                fix_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        fix_text_frame(cell.text_frame)


def main():
    if not os.path.exists(TEMPLATE):
        print(f"找不到模板 {TEMPLATE}，请把你的模板放到 bot/template.pptx", file=sys.stderr)
        sys.exit(1)

    d = this_thursday()
    date_str = d.strftime("%Y-%m-%d")          # 2026-06-25
    title_str = d.strftime("%Y.%m.%d") + "-checkin"   # 2026.06.25-checkin

    prs = Presentation(TEMPLATE)
    replace_placeholders(prs, {"{{DATE}}": date_str, "{{TITLE}}": title_str})

    os.makedirs("out", exist_ok=True)
    out_path = os.path.join("out", f"{date_str}-checkin.pptx")
    prs.save(out_path)

    print(f"✅ 已生成 {out_path}")
    # 供 GitHub Actions 后续步骤使用
    gh_out = os.environ.get("GITHUB_OUTPUT")
    if gh_out:
        with open(gh_out, "a") as f:
            f.write(f"ppt_path={out_path}\n")
            f.write(f"ppt_date={date_str}\n")
            f.write(f"ppt_title={title_str}\n")


if __name__ == "__main__":
    main()
